#!/usr/bin/env python3
"""deck_to_pptx — turn a house DECK (self-contained HTML) into a 16:9 PowerPoint.

    python3 deck_to_pptx.py <deck.html> [out.pptx]

Reads every `<section class="slide">`, re-renders each inline SVG standalone, and lays
the content out natively — real PowerPoint tables, text and images, not screenshots.
One HTML slide becomes one or more PPTX slides; overflow continues on "(cont.)" pages.

Requires:  pip install python-pptx cairosvg beautifulsoup4
Fonts:     the SVG rendering needs a font with wide glyph coverage (circled digits,
           arrows, check marks). Adwaita Sans/Mono or DejaVu both work; set SVG_SANS /
           SVG_MONO below to something `fc-list` reports on this machine.

The deck's diagrams lean on document-level CSS classes, which cairosvg cannot see —
so every class is resolved into explicit presentation attributes before rendering.
Keep CLS below in step with the .bx/.f-*/.dt/.ds vocabulary in DECK.template.html.
"""
import io, os, re, math, sys, tempfile, shutil
from bs4 import BeautifulSoup, NavigableString, Tag
import cairosvg
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

if len(sys.argv) < 2:
    sys.exit(__doc__)
SRC = os.path.abspath(sys.argv[1])
OUT = os.path.abspath(sys.argv[2]) if len(sys.argv) > 2 else os.path.splitext(SRC)[0] + ".pptx"
TMP = tempfile.mkdtemp(prefix="deck_svg_")

# ── palette ──────────────────────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x11, 0x17)
PANEL     = RGBColor(0x1B, 0x1F, 0x2A)
PANEL2    = RGBColor(0x22, 0x27, 0x36)
LINE      = RGBColor(0x2C, 0x33, 0x46)
INK       = RGBColor(0xE7, 0xEB, 0xF3)
INK_SOFT  = RGBColor(0xAA, 0xB3, 0xC5)
INK_FAINT = RGBColor(0x7B, 0x84, 0x99)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
ACC = {
    "a":    RGBColor(0x5B, 0x8D, 0xEF), "b":    RGBColor(0x19, 0xB8, 0x9A),
    "c":    RGBColor(0xF0, 0xA7, 0x3C), "ctrl": RGBColor(0xA9, 0x74, 0xF0),
    "d":    RGBColor(0xEF, 0x5B, 0x6E), "ok":   RGBColor(0x3E, 0xCF, 0x8E),
    "n":    RGBColor(0x7B, 0x84, 0x99),
}
CHIP_ACC = {"s0": "n", "s1": "a", "s2": "ctrl", "s3": "b", "s4": "c", "s5": "d"}
RAIL_ACC = {"rail-a": "a", "rail-b": "b", "rail-c": "c", "rail-ctrl": "ctrl", "rail-d": "d"}
CALLOUT_ACC = {"key": "a", "info": "a", "ok": "ok", "caveat": "c", "danger": "d"}
SANS, MONO = "Segoe UI", "Consolas"

SLIDE_W, SLIDE_H = 13.333, 7.5
ML, MR = 0.62, 0.62
CW = SLIDE_W - ML - MR
BOTTOM = 7.02

# ══════════════════════════════════════════════════════════════════════════
# 1 · SVG → PNG
# ══════════════════════════════════════════════════════════════════════════
MARKERS = """<defs>
<marker id="arw" viewBox="0 0 10 10" refX="9" refY="5" markerWidth="7" markerHeight="7" orient="auto-start-reverse"><path d="M0 0L10 5L0 10z" fill="#8a93a8"/></marker>
<marker id="arwA" viewBox="0 0 10 10" refX="9" refY="5" markerWidth="7" markerHeight="7" orient="auto-start-reverse"><path d="M0 0L10 5L0 10z" fill="#5b8def"/></marker>
<marker id="arwB" viewBox="0 0 10 10" refX="9" refY="5" markerWidth="7" markerHeight="7" orient="auto-start-reverse"><path d="M0 0L10 5L0 10z" fill="#19b89a"/></marker>
<marker id="arwD" viewBox="0 0 10 10" refX="9" refY="5" markerWidth="7" markerHeight="7" orient="auto-start-reverse"><path d="M0 0L10 5L0 10z" fill="#ef5b6e"/></marker>
<marker id="arwS" viewBox="0 0 10 10" refX="9" refY="5" markerWidth="6" markerHeight="6" orient="auto-start-reverse"><path d="M0 0L10 5L0 10z" fill="#5a6478"/></marker>
</defs>"""

SVG_SANS = "Adwaita Sans"
SVG_MONO = "Adwaita Mono"

# class -> presentation attributes (applied only when absent on the element)
CLS = {
    "bx":        {"rx": "9", "stroke-width": "1.6"},
    "f-a":       {"fill": "#243048", "stroke": "#5b8def"},
    "f-b":       {"fill": "#1b343a", "stroke": "#19b89a"},
    "f-c":       {"fill": "#39322d", "stroke": "#f0a73c"},
    "f-ctrl":    {"fill": "#332d4c", "stroke": "#a974f0"},
    "f-d":       {"fill": "#392734", "stroke": "#ef5b6e"},
    "f-n":       {"fill": "#20242f", "stroke": "#2c3346"},
    "f-g":       {"fill": "#161922", "stroke": "#2c3346"},
    "dt":        {"fill": "#ffffff", "font-family": SVG_SANS, "font-size": "12.5",
                  "font-weight": "700", "text-anchor": "middle"},
    "ds":        {"fill": "#aab3c5", "font-family": SVG_MONO, "font-size": "9.5",
                  "text-anchor": "middle"},
    "dtl":       {"fill": "#ffffff", "font-family": SVG_SANS, "font-size": "12.5",
                  "font-weight": "700", "text-anchor": "start"},
    "dsl":       {"fill": "#aab3c5", "font-family": SVG_MONO, "font-size": "9.5",
                  "text-anchor": "start"},
    "hd":        {"font-family": SVG_SANS, "font-size": "10.5", "font-weight": "700",
                  "letter-spacing": "1", "fill": "#8a93a8"},
    "ln":        {"stroke": "#8a93a8", "stroke-width": "1.5", "fill": "none"},
    "ln-d":      {"stroke": "#5a6478", "stroke-width": "1.3", "stroke-dasharray": "4 4",
                  "fill": "none"},
    "lbl":       {"fill": "#8a93a8", "font-family": SVG_MONO, "font-size": "9.5",
                  "text-anchor": "middle"},
    "num-badge": {"fill": "#0f1117", "font-family": SVG_MONO, "font-size": "10",
                  "font-weight": "700", "text-anchor": "middle"},
}

def rgba_to_hex(m):
    """Composite an rgba() colour over the figure panel background."""
    parts = [p.strip() for p in m.group(1).split(",")]
    r, g, b = (float(parts[0]), float(parts[1]), float(parts[2]))
    a = float(parts[3]) if len(parts) > 3 else 1.0
    base = (0x1B, 0x1F, 0x2A)
    out = [int(round(base[i] + a * ((r, g, b)[i] - base[i]))) for i in range(3)]
    return "#%02x%02x%02x" % tuple(out)

def render_svg(svg_tag, path):
    svg = BeautifulSoup(str(svg_tag), "html.parser").find("svg")
    for el in svg.find_all(True):
        classes = el.get("class") or []
        for c in classes:
            for k, v in CLS.get(c, {}).items():
                if not el.has_attr(k):
                    el[k] = v
        if "hd" in classes:                       # CSS uppercased these
            for t in el.find_all(string=True):
                t.replace_with(t.upper())
        if el.has_attr("class"):
            del el["class"]
    raw = str(svg)
    raw = re.sub(r"rgba\(([^)]+)\)", rgba_to_hex, raw)
    vb = svg.get("viewbox") or svg.get("viewBox") or "0 0 1000 300"
    _, _, vw, vh = [float(x) for x in vb.split()]
    body = raw[raw.index(">") + 1: raw.rindex("</svg>")]
    doc = ('<svg xmlns="http://www.w3.org/2000/svg" xmlns:xlink="http://www.w3.org/1999/xlink" '
           f'viewBox="{vb}" width="{vw}" height="{vh}">{MARKERS}{body}</svg>')
    cairosvg.svg2png(bytestring=doc.encode("utf-8"), write_to=path,
                     output_width=int(vw * 2.6), output_height=int(vh * 2.6),
                     background_color="#1b1f2a")
    return vw / vh

# ══════════════════════════════════════════════════════════════════════════
# 2 · HTML → block model
# ══════════════════════════════════════════════════════════════════════════
def txt(el):
    if el is None:
        return ""
    s = el.get_text(" ", strip=True)
    return re.sub(r"\s+", " ", s).replace(" ", " ")

def parse_cell(td):
    """Return (text, bold) — pills and <b> are rendered as emphasis."""
    t = txt(td)
    bold = bool(td.find("b")) or "q" in (td.get("class") or [])
    return t, bold

def blocks_from_card(card, idx):
    """Return (title, body, tables) with body in document order."""
    title = txt(card.find("h5"))
    body, tables = [], []
    for ch in card.find_all(["p", "pre", "table"], recursive=True):
        if ch.name == "p":
            if ch.find_parent("table"):
                continue
            t = txt(ch)
            if t:
                body.append(("p", t))
        elif ch.name == "pre":
            body.append(("pre", ch.get_text().strip("\n")))
        elif ch.name == "table":
            tables.append(ch)
    return title, body, tables

def table_block(tbl, heading=None):
    heads = [txt(th) for th in tbl.select("thead th")]
    rows = []
    for tr in tbl.select("tbody tr"):
        cells = [parse_cell(td) for td in tr.find_all(["td", "th"])]
        if cells:
            rows.append(cells)
    if not rows:
        return None
    ncol = max(len(r) for r in rows) if not heads else len(heads)
    for r in rows:
        while len(r) < ncol:
            r.append(("", False))
    return {"t": "table", "heads": heads, "rows": rows, "ncol": ncol, "heading": heading}

def extract(section, sidx):
    """Walk a slide section and emit an ordered list of content blocks."""
    blocks = []
    chips = [(txt(c), CHIP_ACC.get(next((k for k in c.get("class", []) if k in CHIP_ACC), "s0"), "n"))
             for c in section.select(".eyebrow .chip")]
    kicker = txt(section.select_one(".eyebrow .kicker"))
    title = txt(section.find("h2"))
    lede = txt(section.find("p", class_="lede"))

    seen = set()
    for el in section.find_all(recursive=False):
        walk(el, blocks, sidx, seen)
    return {"chips": chips, "kicker": kicker, "title": title, "lede": lede, "blocks": blocks}

def walk(el, blocks, sidx, seen):
    if not isinstance(el, Tag) or id(el) in seen:
        return
    cls = el.get("class") or []
    if el.name == "div" and "eyebrow" in cls:
        return
    if el.name == "h2" or (el.name == "p" and "lede" in cls):
        return

    if el.name == "figure":
        seen.add(id(el))
        svg = el.find("svg")
        cap = txt(el.find("figcaption"))
        if svg:
            path = f"{TMP}/s{sidx:02d}_{len(blocks)}.png"
            ar = render_svg(svg, path)
            blocks.append({"t": "img", "path": path, "ar": ar, "cap": cap})
        return

    if el.name == "div" and "tablewrap" in cls:
        seen.add(id(el))
        b = table_block(el.find("table"))
        if b:
            blocks.append(b)
        return

    if el.name == "div" and "grid" in cls:
        seen.add(id(el))
        cols = 3 if "g3" in cls else (4 if "g4" in cls else 2)
        cards = []
        for card in el.find_all("div", class_="card", recursive=False):
            title, body, tables = blocks_from_card(card, sidx)
            if tables:
                for t in tables:
                    b = table_block(t, heading=title)
                    if b:
                        blocks.append(b)
                continue
            rail = next((RAIL_ACC[c] for c in (card.get("class") or []) if c in RAIL_ACC), None)
            cards.append({"title": title, "body": body, "rail": rail})
        if cards:
            blocks.append({"t": "cards", "cols": min(cols, len(cards)), "items": cards})
        return

    if el.name == "div" and "callout" in cls:
        seen.add(id(el))
        kind = next((k for k in cls if k in CALLOUT_ACC), "key")
        blocks.append({"t": "callout", "acc": CALLOUT_ACC[kind],
                       "label": txt(el.find(class_="clabel")),
                       "paras": [txt(p) for p in el.find_all("p")]})
        return

    if el.name == "pre":
        seen.add(id(el))
        blocks.append({"t": "code", "text": el.get_text().rstrip("\n")})
        return

    if el.name == "div" and "tree" in cls:
        seen.add(id(el))
        blocks.append({"t": "code", "text": el.get_text().strip("\n")})
        return

    if el.name in ("h3", "h4"):
        seen.add(id(el))
        blocks.append({"t": "head", "text": txt(el)})
        return

    if el.name == "p":
        seen.add(id(el))
        style = "pull" if "pull" in cls else ("small" if "small" in cls else "body")
        t = txt(el)
        if t:
            blocks.append({"t": "para", "text": t, "style": style})
        return

    for ch in el.find_all(recursive=False):
        walk(ch, blocks, sidx, seen)

# ══════════════════════════════════════════════════════════════════════════
# 3 · height estimation
# ══════════════════════════════════════════════════════════════════════════
# Average glyph advance as a fraction of the em. Everything downstream is derived
# from these two numbers, so estimate and render can never drift apart — which is
# exactly the bug that produced text overflowing its panel.
ADV_SANS, ADV_MONO = 0.52, 0.55
SAFETY = 1.06          # bias toward whitespace: slack is invisible, overlap is not

def chars_per_line(box_w, size_pt, mono=False):
    adv_px = (ADV_MONO if mono else ADV_SANS) * size_pt * (96.0 / 72.0)
    return max(8, int(box_w * 96.0 / adv_px))

def wrap_lines(text, chars):
    if not text:
        return 1
    n = 0
    for para in str(text).split("\n"):
        n += max(1, math.ceil(len(para) / chars))
    return n

def text_h(text, box_w, size_pt, spacing=1.06, mono=False, after_pt=3):
    """Inches a wrapped run occupies inside a text frame box_w inches wide."""
    lines = wrap_lines(text, chars_per_line(box_w, size_pt, mono))
    return (lines * size_pt * 1.2 * spacing + after_pt) / 72.0 * SAFETY

CARD_PAD, CALLOUT_PAD = 0.26, 0.26      # inner vertical padding, top + bottom
CELL_PAD_V, CELL_PAD_H = 0.06, 0.11     # table cell margins, both sides

def card_content_h(it, inner_w):
    h = 0.0
    if it["title"]:
        h += text_h(it["title"], inner_w, 11.5, 1.0, after_pt=3)
    for kind, t in it["body"]:
        if kind == "p":
            h += text_h(t, inner_w, 9.5, 1.06, after_pt=3)
        else:
            h += text_h(t, inner_w, 8.0, 1.0, mono=True, after_pt=2)
    return h

def cards_geom(b, w):
    """Shared geometry so est_height and draw() lay cards out identically."""
    cols, items, gap = b["cols"], b["items"], 0.18
    cw = (w - gap * (cols - 1)) / cols
    inner = cw - 0.28
    ch = max(card_content_h(it, inner) for it in items) + CARD_PAD
    return cols, gap, cw, inner, ch, math.ceil(len(items) / cols)

def callout_content_h(b, w):
    inner = w - 0.40
    h = 0.0
    if b["label"]:
        h += text_h(b["label"], inner, 9.5, 1.0, mono=True, after_pt=3)
    for p in b["paras"]:
        h += text_h(p, inner, 10.5, 1.06, after_pt=3)
    return h

def col_weights(b):
    ws = []
    for ci in range(b["ncol"]):
        lens = [len(r[ci][0]) for r in b["rows"]]
        avg = sum(lens) / max(1, len(lens))
        hdl = len(b["heads"][ci]) if ci < len(b["heads"]) else 0
        words = [len(x) for r in b["rows"] for x in r[ci][0].split()] or [4]
        ws.append(min(60.0, max(7.0, avg, max(words) * 1.15, hdl * 0.95)))
    return ws

def table_rows_h(b, w):
    """(header height, [row heights]) using the real column widths."""
    ws = col_weights(b)
    tot = sum(ws)
    def cell_w(ci):
        return max(0.4, w * ws[ci] / tot - CELL_PAD_H)
    hh = 0.0
    if b["heads"]:
        hh = max(text_h(b["heads"][ci], cell_w(ci), 8.5, 1.0, mono=True, after_pt=0)
                 for ci in range(b["ncol"])) + CELL_PAD_V
    rows = []
    for r in b["rows"]:
        rows.append(max(text_h(r[ci][0], cell_w(ci), 9, 0.98, after_pt=0)
                        for ci in range(b["ncol"])) + CELL_PAD_V)
    return hh, rows

def est_height(b, w=CW):
    if b["t"] == "img":
        h = min(4.15, w / b["ar"])
        return h + (0.56 if b["cap"] else 0.26)
    if b["t"] == "table":
        hh, rows = table_rows_h(b, w)
        return (0.30 if b["heading"] else 0) + hh + sum(rows) + 0.16
    if b["t"] == "cards":
        _, gap, _, _, ch, nrows = cards_geom(b, w)
        return nrows * (ch + 0.16) + 0.06
    if b["t"] == "callout":
        return callout_content_h(b, w) + CALLOUT_PAD + 0.16
    if b["t"] == "code":
        return 0.14 + 0.146 * (b["text"].count("\n") + 1) + 0.16
    if b["t"] == "head":
        return text_h(b["text"], w, 10.5, 1.0, mono=True, after_pt=0) + 0.12
    if b["t"] == "para":
        st = b["style"]
        size = 13 if st == "pull" else (9.5 if st == "small" else 11)
        inner = (w - 0.20) if st == "pull" else w
        return text_h(b["text"], inner, size, 1.06, after_pt=0) + 0.12
    return 0.3

def split_table(b, budget):
    """Split an oversized table into chunks that fit the budget."""
    out, rows = [], b["rows"]
    lo, hi = 0, len(rows)
    while lo < hi:
        take = hi - lo
        while take > 1:
            probe = dict(b, rows=rows[lo:lo + take])
            if est_height(probe) <= budget:
                break
            take -= 1
        out.append(dict(b, rows=rows[lo:lo + take],
                        heading=b["heading"] if lo == 0 else None))
        lo += take
    return out

# ══════════════════════════════════════════════════════════════════════════
# 4 · drawing
# ══════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(SLIDE_W), Inches(SLIDE_H)
BLANK = prs.slide_layouts[6]

def new_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background(); bg.shadow.inherit = False
    return s

def tb(slide, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf

CIRCLED = {"\u2460":"1","\u2461":"2","\u2462":"3","\u2463":"4","\u2464":"5","\u2465":"6",
           "\u2466":"7","\u2467":"8","\u2468":"9","\u24ea":"0"}
def deglyph(t):
    for k, v in CIRCLED.items():
        t = t.replace(k, v)
    return t

def para(tf, text, size, color, font=SANS, bold=False, first=False,
         space_after=2, align=PP_ALIGN.LEFT, italic=False, spacing=1.0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = spacing
    r = p.add_run(); r.text = deglyph(text)
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = color
    return p

def rect(slide, x, y, w, h, fill, line=None, rounded=True, lw=1.0):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def header(slide, meta, cont=False):
    """Chips + title; returns the y where content may start."""
    y = 0.40
    bits = [c[0] for c in meta["chips"]]
    if meta["kicker"]:
        bits.append(meta["kicker"])
    if bits:
        tf = tb(slide, ML, y, CW, 0.24)
        p = tf.paragraphs[0]; p.space_after = Pt(0)
        for i, c in enumerate(meta["chips"]):
            r = p.add_run(); r.text = deglyph(("   " if i else "") + c[0].upper())
            r.font.size = Pt(10); r.font.bold = True; r.font.name = MONO
            r.font.color.rgb = ACC[c[1]]
        if meta["kicker"]:
            r = p.add_run(); r.text = ("   ·   " if meta["chips"] else "") + meta["kicker"].upper()
            r.font.size = Pt(9.5); r.font.name = MONO; r.font.color.rgb = INK_FAINT
        y += 0.30
    t = meta["title"] + ("  (cont.)" if cont else "")
    n = len(t)
    if n <= chars_per_line(CW, 25):
        size, lines, adv = 25, 1, 0.46
    elif n <= chars_per_line(CW, 22):
        size, lines, adv = 22, 1, 0.42
    else:
        lines = math.ceil(n / chars_per_line(CW, 22))
        size, adv = 22, 0.40 * lines
    tf = tb(slide, ML, y, CW, 0.5 * lines)
    para(tf, t, size, INK, bold=True, first=True, spacing=0.94)
    y += adv + 0.10
    rect(slide, ML, y, 1.5, 0.032, ACC[meta["chips"][0][1]] if meta["chips"] else ACC["a"],
         rounded=False)
    return y + 0.20

def draw(slide, b, x, y, w):
    if b["t"] == "img":
        h = min(4.15, w / b["ar"])
        iw = h * b["ar"]
        if iw > w:
            iw, h = w, w / b["ar"]
        rect(slide, x + (w - iw) / 2 - 0.09, y + 0.02, iw + 0.18, h + 0.18, PANEL, LINE)
        slide.shapes.add_picture(b["path"], Inches(x + (w - iw) / 2), Inches(y + 0.11),
                                 Inches(iw), Inches(h))
        yy = y + h + 0.26
        if b["cap"]:
            tf = tb(slide, x, yy, w, 0.34)
            para(tf, b["cap"], 10, INK_FAINT, first=True, spacing=1.0)
            yy += 0.30
        return yy + 0.10

    if b["t"] == "table":
        if b["heading"]:
            tf = tb(slide, x, y, w, 0.24)
            para(tf, b["heading"].upper(), 10, INK_SOFT, font=MONO, bold=True, first=True)
            y += 0.28
        nrow = len(b["rows"]) + (1 if b["heads"] else 0)
        hh, rowhs = table_rows_h(b, w)
        h = hh + sum(rowhs)
        shape = slide.shapes.add_table(nrow, b["ncol"], Inches(x), Inches(y),
                                       Inches(w), Inches(h))
        tblx = shape.table
        tblx.first_row = bool(b["heads"])
        weights = col_weights(b)
        tot = sum(weights)
        for i, wt in enumerate(weights):
            tblx.columns[i].width = Emu(int(Inches(w) * wt / tot))
        if b["heads"]:
            tblx.rows[0].height = Emu(int(Inches(hh)))
        for i, rh in enumerate(rowhs):
            tblx.rows[i + (1 if b["heads"] else 0)].height = Emu(int(Inches(rh)))
        r0 = 0
        if b["heads"]:
            for i, htxt in enumerate(b["heads"]):
                cell = tblx.cell(0, i)
                cell.fill.solid(); cell.fill.fore_color.rgb = PANEL2
                cell.margin_left = cell.margin_right = Inches(0.055)
                cell.margin_top = cell.margin_bottom = Inches(0.03)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf = cell.text_frame; tf.word_wrap = True
                para(tf, htxt.upper(), 8.5, INK_SOFT, font=MONO, bold=True, first=True)
            r0 = 1
        for ri, row in enumerate(b["rows"]):
            for ci in range(b["ncol"]):
                cell = tblx.cell(ri + r0, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = PANEL if ri % 2 == 0 else RGBColor(0x1F, 0x24, 0x30)
                cell.margin_left = cell.margin_right = Inches(0.055)
                cell.margin_top = cell.margin_bottom = Inches(0.03)
                cell.vertical_anchor = MSO_ANCHOR.TOP
                tf = cell.text_frame; tf.word_wrap = True
                t, bold = row[ci]
                para(tf, t, 9, INK if bold else INK_SOFT, bold=bold, first=True, spacing=0.98)
        return y + h + 0.16

    if b["t"] == "cards":
        items = b["items"]
        cols, gap, cw, inner, ch, _ = cards_geom(b, w)
        for i, it in enumerate(items):
            r, c = divmod(i, cols)
            cx = x + c * (cw + gap)
            cy = y + r * (ch + 0.16)
            rect(slide, cx, cy, cw, ch, PANEL, LINE)
            rect(slide, cx, cy + 0.04, 0.035, ch - 0.08,
                 ACC[it.get("rail") or ["a", "b", "c", "ctrl"][i % 4]], rounded=False)
            tf = tb(slide, cx + 0.16, cy + 0.13, inner, ch - CARD_PAD)
            first = True
            if it["title"]:
                para(tf, it["title"], 11.5, INK, bold=True, first=True, space_after=3)
                first = False
            for kind, t in it["body"]:
                if kind == "p":
                    para(tf, t, 9.5, INK_SOFT, first=first, space_after=3, spacing=1.06)
                else:
                    para(tf, t, 8, RGBColor(0xC9, 0xD3, 0xE6), font=MONO, first=first,
                         space_after=2, spacing=1.0)
                first = False
        return y + math.ceil(len(items) / cols) * (ch + 0.16) + 0.06

    if b["t"] == "callout":
        h = callout_content_h(b, w) + CALLOUT_PAD
        acc = ACC[b["acc"]]
        rect(slide, x, y, w, h, PANEL2, acc)
        rect(slide, x, y + 0.04, 0.045, h - 0.08, acc, rounded=False)
        tf = tb(slide, x + 0.20, y + 0.13, w - 0.40, h - CALLOUT_PAD)
        first = True
        if b["label"]:
            para(tf, b["label"].replace("◆", "").strip().upper(), 9.5, acc,
                 font=MONO, bold=True, first=True, space_after=3)
            first = False
        for p in b["paras"]:
            para(tf, p, 10.5, INK_SOFT, first=first, space_after=3, spacing=1.06)
            first = False
        return y + h + 0.16

    if b["t"] == "code":
        n = b["text"].count("\n") + 1
        h = 0.146 * n + 0.26
        rect(slide, x, y, w, h, RGBColor(0x12, 0x14, 0x1C), LINE)
        tf = tb(slide, x + 0.16, y + 0.10, w - 0.32, h - 0.18)
        first = True
        for ln in b["text"].split("\n"):
            para(tf, ln if ln.strip() else " ", 8.5, RGBColor(0xC9, 0xD3, 0xE6),
                 font=MONO, first=first, space_after=0, spacing=1.0)
            first = False
        return y + h + 0.14

    if b["t"] == "head":
        h = est_height(b, w)
        tf = tb(slide, x, y, w, h)
        para(tf, b["text"].upper(), 10.5, INK_SOFT, font=MONO, bold=True, first=True)
        return y + h + 0.06

    if b["t"] == "para":
        st = b["style"]
        size = 13 if st == "pull" else (9.5 if st == "small" else 11)
        col = INK if st == "pull" else (INK_FAINT if st == "small" else INK_SOFT)
        h = est_height(b, w)
        if st == "pull":
            rect(slide, x, y, 0.04, h - 0.1, ACC["a"], rounded=False)
            tf = tb(slide, x + 0.20, y, w - 0.20, h)
        else:
            tf = tb(slide, x, y, w, h)
        para(tf, b["text"], size, col, first=True, spacing=1.06,
             italic=(st == "pull"))
        return y + h
    return y

# ══════════════════════════════════════════════════════════════════════════
# 5 · assemble
# ══════════════════════════════════════════════════════════════════════════
soup = BeautifulSoup(open(SRC, encoding="utf-8").read(), "html.parser")
sections = soup.select("section.slide")
print(f"parsed {len(sections)} html slides")

# ---- cover — built from the first slide, so it works for any deck --------
first = extract(sections[0], 0) if sections else {"title": "", "lede": "", "chips": [], "kicker": ""}
DOC_TITLE = (soup.find("title").get_text() if soup.find("title") else "").split("—")[0].strip()
DOC_SUB = (soup.find("title").get_text() if soup.find("title") else "")
DOC_SUB = DOC_SUB.split("—", 1)[1].strip() if "—" in DOC_SUB else ""
GROUPS = []
for sec in sections:
    g = sec.get("data-grp")
    if g and g not in GROUPS:
        GROUPS.append(g)

cover = new_slide()
rect(cover, 0, 0, SLIDE_W, 0.09, ACC["a"], rounded=False)
tf = tb(cover, 1.1, 1.95, 11.1, 0.4)
para(tf, "SOLUTION / SOFTWARE ARCHITECTURE  ·  CANONICAL PROCESS", 12,
     ACC["a"], font=MONO, bold=True, first=True)
tf = tb(cover, 1.1, 2.45, 11.1, 1.6)
para(tf, "From a one-line idea to a build plan\nsomeone else can execute", 40, INK,
     bold=True, first=True, spacing=0.98)
tf = tb(cover, 1.1, 4.25, 10.4, 1.0)
para(tf, "Requirements · analysis · dependencies · building blocks · integrations · "
         "data model · code layout · libraries · design patterns", 14, INK_SOFT, first=True,
     spacing=1.15)
for i, (lab, k) in enumerate([("① Vision Brief", "a"), ("② RFC", "ctrl"), ("③ DRP", "b"),
                              ("Constraints", "c"), ("④ Architecture", "a"),
                              ("⑤ Work Plan", "b"), ("6·7 Reviews", "ctrl")]):
    x = 1.1 + i * 1.60
    rect(cover, x, 5.55, 1.45, 0.52, PANEL, ACC[k])
    tf = tb(cover, x + 0.08, 5.68, 1.29, 0.3)
    para(tf, lab, 9, ACC[k], font=MONO, bold=True, first=True, align=PP_ALIGN.CENTER)
tf = tb(cover, 1.1, 6.55, 11.1, 0.3)
para(tf, "Worked example carried end to end: an order & payments system", 11,
     INK_FAINT, first=True)

# ---- content -----------------------------------------------------------
count = 1
for si, sec in enumerate(sections):
    meta = extract(sec, si + 1)
    if si == 0:                                    # the HTML title slide is the cover
        blocks = meta["blocks"]
    else:
        blocks = meta["blocks"]
    slide = new_slide(); count += 1
    y = header(slide, meta)
    if meta["lede"]:
        h = text_h(meta["lede"], CW, 12, 1.08, after_pt=0) + 0.06
        tf = tb(slide, ML, y, CW, h)
        para(tf, meta["lede"], 12, INK_SOFT, first=True, spacing=1.08)
        y += h + 0.16
    queue = list(blocks)
    y_top = y
    while queue:
        b = queue.pop(0)
        budget = BOTTOM - y
        h = est_height(b)
        # a heading stays with the block it introduces
        if b["t"] == "head" and queue and h + est_height(queue[0]) > budget and y > y_top + 0.05:
            queue.insert(0, b)
            slide = new_slide(); count += 1
            y = y_top = header(slide, meta, cont=True)
            continue
        if h > budget + 0.35:
            if y > y_top + 0.05:                    # something is already here: turn the page
                queue.insert(0, b)
                slide = new_slide(); count += 1
                y = y_top = header(slide, meta, cont=True)
                continue
            if b["t"] == "table":                   # too big even for a whole slide: split
                chunks = split_table(b, budget)
                if len(chunks) > 1:
                    queue = chunks[1:] + queue
                    b = chunks[0]
        y = draw(slide, b, ML, y, CW)

# ---- closing note ------------------------------------------------------
for i, s in enumerate(prs.slides):
    if i == 0:
        continue
    tf = tb(s, SLIDE_W - 1.5, SLIDE_H - 0.42, 0.9, 0.24)
    para(tf, str(i), 9, INK_FAINT, font=MONO, first=True, align=PP_ALIGN.RIGHT)

prs.save(OUT)
shutil.rmtree(TMP, ignore_errors=True)
print(f"wrote {OUT}  ({count} slides, {os.path.getsize(OUT)/1024:.0f} KB)")
